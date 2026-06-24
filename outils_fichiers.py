"""
NEOGEN - Outils de lecture et création de fichiers.
Formats : PDF, PPTX, DOCX (lecture + création DOCX).
Dépendances optionnelles : pypdf, python-pptx, python-docx
"""
from __future__ import annotations
import base64
import io
import os
import uuid

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
RAPPORTS_DIR = os.path.join(BASE_DIR, "data", "rapports")
MAX_CHARS = 12000


def extraire_texte_b64(fichier_b64: str, nom: str) -> str:
    """Extrait le texte d'un fichier encodé en base64 (PDF/PPTX/DOCX/TXT)."""
    ext = nom.lower().rsplit(".", 1)[-1] if "." in nom else ""
    data = base64.b64decode(fichier_b64)
    return _extraire(data, ext, nom)


def lire_fichier_chemin(chemin: str) -> str:
    """Lit un fichier depuis un chemin local (PDF/PPTX/DOCX/TXT)."""
    if not os.path.exists(chemin):
        return f"[Fichier introuvable : {chemin}]"
    ext = chemin.lower().rsplit(".", 1)[-1] if "." in chemin else ""
    with open(chemin, "rb") as f:
        data = f.read()
    return _extraire(data, ext, os.path.basename(chemin))


def _extraire(data: bytes, ext: str, nom: str) -> str:
    if ext == "pdf":
        return _lire_pdf(data)
    if ext in ("pptx", "ppt"):
        return _lire_pptx(data)
    if ext in ("docx", "doc"):
        return _lire_docx(data)
    if ext in ("txt", "md", "csv", "html", "htm", "json", "xml", "log"):
        try:
            return data.decode("utf-8", errors="replace")[:MAX_CHARS]
        except Exception:
            return "[Lecture texte échouée]"
    return f"[Format .{ext} non supporté — acceptés : PDF, PPTX, DOCX, TXT, HTML, CSV, JSON]"


def _lire_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages = []
        for i, page in enumerate(reader.pages, 1):
            t = (page.extract_text() or "").strip()
            if t:
                pages.append(f"[Page {i}]\n{t}")
        return "\n\n".join(pages)[:MAX_CHARS] or "[PDF sans texte extractible]"
    except ImportError:
        return "[pypdf non installé — rebuild Docker requis]"
    except Exception as e:
        return f"[Erreur lecture PDF : {e}]"


def _lire_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            textes = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
            if textes:
                slides.append(f"[Diapo {i}]\n" + "\n".join(textes))
        return "\n\n".join(slides)[:MAX_CHARS] or "[PPTX sans texte extractible]"
    except ImportError:
        return "[python-pptx non installé — rebuild Docker requis]"
    except Exception as e:
        return f"[Erreur lecture PPTX : {e}]"


def _lire_docx(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paras)[:MAX_CHARS] or "[DOCX sans texte extractible]"
    except ImportError:
        return "[python-docx non installé — rebuild Docker requis]"
    except Exception as e:
        return f"[Erreur lecture DOCX : {e}]"


def creer_rapport_pdf(titre: str, contenu: str) -> str | None:
    """Crée un rapport PDF structuré. Retourne le nom de fichier ou None si fpdf2 absent."""
    try:
        from fpdf import FPDF
    except ImportError:
        return None
    os.makedirs(RAPPORTS_DIR, exist_ok=True)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Helvetica", "B", 16)
    pdf.multi_cell(0, 10, (titre or "Rapport NEOGEN").encode("latin-1", "replace").decode("latin-1"))
    pdf.ln(4)
    pdf.set_font("Helvetica", size=11)
    for ligne in contenu.splitlines():
        if ligne.startswith("## ") or ligne.startswith("### "):
            pdf.set_font("Helvetica", "B", 12)
            texte = ligne.lstrip("#").strip()
            pdf.multi_cell(0, 8, texte.encode("latin-1", "replace").decode("latin-1"))
            pdf.set_font("Helvetica", size=11)
        elif ligne.strip():
            pdf.multi_cell(0, 7, ligne.encode("latin-1", "replace").decode("latin-1"))
        else:
            pdf.ln(3)
    nom = f"rapport_{uuid.uuid4().hex[:8]}.pdf"
    pdf.output(os.path.join(RAPPORTS_DIR, nom))
    return nom


def creer_rapport_excel(titre: str, contenu: str) -> str | None:
    """Crée un fichier Excel à partir de contenu tabulaire (CSV ou lignes séparées par |).
    Retourne le nom de fichier ou None si openpyxl absent."""
    try:
        from openpyxl import Workbook
    except ImportError:
        return None
    os.makedirs(RAPPORTS_DIR, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = (titre or "Données")[:31]
    lignes = [l for l in contenu.splitlines() if l.strip()]
    for row_idx, ligne in enumerate(lignes, 1):
        if "|" in ligne:
            cellules = [c.strip() for c in ligne.split("|") if c.strip()]
        elif "," in ligne:
            import csv as _csv
            import io as _io
            cellules = next(_csv.reader(_io.StringIO(ligne)))
        else:
            cellules = [ligne.strip()]
        for col_idx, val in enumerate(cellules, 1):
            ws.cell(row=row_idx, column=col_idx, value=val)
    nom = f"rapport_{uuid.uuid4().hex[:8]}.xlsx"
    wb.save(os.path.join(RAPPORTS_DIR, nom))
    return nom


def creer_rapport_csv(titre: str, contenu: str) -> str | None:
    """Crée un fichier CSV. Retourne le nom de fichier."""
    os.makedirs(RAPPORTS_DIR, exist_ok=True)
    nom = f"rapport_{uuid.uuid4().hex[:8]}.csv"
    with open(os.path.join(RAPPORTS_DIR, nom), "w", encoding="utf-8-sig", newline="") as f:
        if titre:
            f.write(f"# {titre}\n")
        f.write(contenu)
    return nom


def creer_rapport_pptx(titre: str, contenu: str) -> str | None:
    """Crée une présentation PowerPoint. Retourne le nom de fichier ou None si python-pptx absent."""
    try:
        from pptx import Presentation as _Prs
    except ImportError:
        return None
    os.makedirs(RAPPORTS_DIR, exist_ok=True)
    prs = _Prs()
    # Diapo titre
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = titre or "Présentation NEOGEN"
    if slide.placeholders[1]:
        slide.placeholders[1].text = "Généré par NEOGEN"
    # Diapos contenu : chaque section ## / ### devient une diapo
    slide_titre = None
    corps: list[str] = []

    def _flush():
        nonlocal slide_titre, corps
        if slide_titre is None and not corps:
            return
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = slide_titre or "Contenu"
        tf = s.placeholders[1].text_frame
        tf.text = "\n".join(corps) if corps else ""
        slide_titre, corps = None, []

    for ligne in contenu.splitlines():
        if ligne.startswith("## ") or ligne.startswith("### "):
            _flush()
            slide_titre = ligne.lstrip("#").strip()
        elif ligne.strip():
            if slide_titre is None:
                slide_titre = "Contenu"
            corps.append(ligne.strip())
    _flush()
    nom = f"rapport_{uuid.uuid4().hex[:8]}.pptx"
    prs.save(os.path.join(RAPPORTS_DIR, nom))
    return nom


def creer_rapport_html(titre: str, contenu: str) -> str | None:
    """Crée un rapport HTML formaté. Retourne le nom de fichier."""
    import html as _html
    os.makedirs(RAPPORTS_DIR, exist_ok=True)
    lignes_html = []
    for ligne in contenu.splitlines():
        if ligne.startswith("### "):
            lignes_html.append(f"<h3>{_html.escape(ligne[4:])}</h3>")
        elif ligne.startswith("## "):
            lignes_html.append(f"<h2>{_html.escape(ligne[3:])}</h2>")
        elif ligne.startswith("# "):
            lignes_html.append(f"<h2>{_html.escape(ligne[2:])}</h2>")
        elif ligne.strip():
            lignes_html.append(f"<p>{_html.escape(ligne)}</p>")
        else:
            lignes_html.append("<br>")
    titre_esc = _html.escape(titre or "Rapport NEOGEN")
    html_doc = (
        f'<!DOCTYPE html><html lang="fr"><head><meta charset="UTF-8"><title>{titre_esc}</title>'
        '<style>body{font-family:Arial,sans-serif;max-width:820px;margin:40px auto;padding:0 24px;color:#333}'
        'h1{color:#1a1a2e}h2{color:#16213e;border-bottom:1px solid #eee;padding-bottom:6px}'
        'h3{color:#0f3460}p{line-height:1.6}</style></head>'
        f'<body><h1>{titre_esc}</h1>\n'
        + "\n".join(lignes_html)
        + '\n</body></html>'
    )
    nom = f"rapport_{uuid.uuid4().hex[:8]}.html"
    with open(os.path.join(RAPPORTS_DIR, nom), "w", encoding="utf-8") as f:
        f.write(html_doc)
    return nom


def creer_rapport_docx(titre: str, contenu: str) -> str | None:
    """Crée un rapport DOCX structuré. Retourne le nom de fichier ou None si erreur."""
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        return None
    os.makedirs(RAPPORTS_DIR, exist_ok=True)
    doc = Document()
    doc.add_heading(titre or "Rapport NEOGEN", level=1)
    for ligne in contenu.splitlines():
        if ligne.startswith("### "):
            doc.add_heading(ligne[4:], level=3)
        elif ligne.startswith("## "):
            doc.add_heading(ligne[3:], level=2)
        elif ligne.startswith("# "):
            doc.add_heading(ligne[2:], level=1)
        elif ligne.strip():
            doc.add_paragraph(ligne)
    nom = f"rapport_{uuid.uuid4().hex[:8]}.docx"
    doc.save(os.path.join(RAPPORTS_DIR, nom))
    return nom
